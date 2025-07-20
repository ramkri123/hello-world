#!/usr/bin/env python3
"""
Extract all slides from PowerPoint to see what's available
"""

from pptx import Presentation
import os

def extract_all_slides(pptx_path):
    """Extract all slides from PowerPoint"""
    
    try:
        # Load presentation
        print(f"📖 Loading PowerPoint: {pptx_path}")
        prs = Presentation(pptx_path)
        
        print(f"📊 Total slides found: {len(prs.slides)}")
        
        markdown_content = f"# PowerPoint Content Extract - All Slides\n"
        markdown_content += f"**Source:** {os.path.basename(pptx_path)}\n"
        markdown_content += f"**Total Slides:** {len(prs.slides)}\n\n"
        
        # Extract all slides
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            print(f"📄 Extracting slide {slide_number}")
            
            markdown_content += f"\n## Slide {slide_number}\n\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    markdown_content += f"{text}\n\n"
        
        # Write to file
        output_path = "all_slides_extracted.md"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        print(f"✅ All slides saved to: {output_path}")
        return True
        
    except Exception as e:
        print(f"❌ Error: {e}")
        return False

if __name__ == "__main__":
    extract_all_slides("Zero-trust Sovereign AI.pptx")
