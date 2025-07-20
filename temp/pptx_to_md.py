#!/usr/bin/env python3
"""
Convert PowerPoint (PPTX) to Markdown format
Extract slides 27, 28, 29 specifically
"""

from pptx import Presentation
import os

def extract_slide_content(slide, slide_number):
    """Extract content from a single slide"""
    content = f"\n## Slide {slide_number}\n\n"
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            
            # Check if it's likely a title (larger text, first shape, etc.)
            if shape.top < 1000000:  # Rough heuristic for title position
                if len(text) < 100 and '\n' not in text:  # Likely a title
                    content += f"### {text}\n\n"
                else:
                    content += f"{text}\n\n"
            else:
                content += f"{text}\n\n"
    
    return content

def convert_pptx_to_md(pptx_path, output_path=None, specific_slides=None):
    """Convert PowerPoint to Markdown"""
    
    if not os.path.exists(pptx_path):
        print(f"❌ File not found: {pptx_path}")
        return False
    
    try:
        # Load presentation
        print(f"📖 Loading PowerPoint: {pptx_path}")
        prs = Presentation(pptx_path)
        
        # Prepare output
        if output_path is None:
            output_path = pptx_path.replace('.pptx', '_extracted.md')
        
        markdown_content = f"# PowerPoint Content Extract\n"
        markdown_content += f"**Source:** {os.path.basename(pptx_path)}\n"
        markdown_content += f"**Total Slides:** {len(prs.slides)}\n\n"
        
        # Extract slides
        slides_to_extract = specific_slides if specific_slides else range(len(prs.slides))
        
        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            
            # Skip if not in specific slides list
            if specific_slides and slide_number not in specific_slides:
                continue
            
            print(f"📄 Extracting slide {slide_number}")
            slide_content = extract_slide_content(slide, slide_number)
            markdown_content += slide_content
        
        # Write to file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        print(f"✅ Markdown saved to: {output_path}")
        return True
        
    except Exception as e:
        print(f"❌ Error converting PowerPoint: {e}")
        return False

def main():
    """Main function to extract specific slides"""
    pptx_file = "Zero-trust Sovereign AI.pptx"
    
    # Extract slides 27, 28, 29 specifically
    target_slides = [27, 28, 29]
    
    print("🚀 PowerPoint to Markdown Converter")
    print(f"Target slides: {target_slides}")
    
    success = convert_pptx_to_md(
        pptx_path=pptx_file,
        output_path="slides_27_28_29_architecture.md",
        specific_slides=target_slides
    )
    
    if success:
        print("\n📋 Next steps:")
        print("1. Review the extracted content")
        print("2. Generalize the architecture for multiple use cases")
        print("3. Implement distributed HTTP-based architecture")
    else:
        print("\n❌ Extraction failed. Please check the file and try again.")

if __name__ == "__main__":
    main()
